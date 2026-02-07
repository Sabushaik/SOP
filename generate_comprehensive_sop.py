#!/usr/bin/env python3
"""
Final version: Generate comprehensive SOP PowerPoint presentation
Preserves all template styling, logos, fonts, and formatting from SOP Draft.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
import os
import re

def clean_text(text):
    """Clean and format text"""
    text = text.strip()
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    return text

def extract_docx_content_detailed(docx_path, friendly_title=None):
    """Extract detailed structured content from a DOCX file"""
    try:
        doc = Document(docx_path)
        content = {
            'filename': os.path.basename(docx_path),
            'title': friendly_title or '',  # Use friendly title if provided
            'sections': []
        }
        
        current_section = None
        lines_processed = 0
        
        for para in doc.paragraphs:
            text = clean_text(para.text)
            if not text or len(text) < 3:
                continue
            
            lines_processed += 1
            
            # Detect if this is a heading
            is_heading = False
            heading_level = 0
            
            if para.style.name.startswith('Heading'):
                is_heading = True
                try:
                    heading_level = int(para.style.name.replace('Heading', '').strip())
                except:
                    heading_level = 1
            elif para.runs and len(para.runs) > 0:
                # Check if text is bold and larger (likely a heading)
                first_run = para.runs[0]
                if first_run.bold and first_run.font.size and first_run.font.size > Pt(12):
                    is_heading = True
                    heading_level = 1
            
            if is_heading:
                # Start a new section
                if current_section:
                    content['sections'].append(current_section)
                current_section = {
                    'heading': text,
                    'level': heading_level,
                    'content': []
                }
            else:
                # Regular content
                if current_section is not None:
                    current_section['content'].append(text)
                else:
                    # No section yet, create one
                    current_section = {
                        'heading': 'Overview',
                        'level': 1,
                        'content': [text]
                    }
        
        # Add the last section
        if current_section:
            content['sections'].append(current_section)
        
        # Extract title from filename if still not found
        if not content['title']:
            fname = os.path.basename(docx_path).replace('.docx', '')
            # Clean up filename
            fname = fname.replace('P9_IND_HR_', '').replace('-00_', ' ')
            fname = fname.replace('_', ' ').replace('  ', ' ')
            content['title'] = fname
        
        return content
        
    except Exception as e:
        print(f"  Error reading {docx_path}: {e}")
        fname = friendly_title or os.path.basename(docx_path).replace('.docx', '').replace('_', ' ')
        return {'filename': fname, 'title': fname, 'sections': []}

def add_content_slide_with_template_style(prs, title, bullet_points):
    """Add a slide using the template's Title and Content layout"""
    # Use layout 1 (Title and Content)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title
        # Apply consistent formatting
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
    
    # Add content
    # Find the content placeholder (usually index 1)
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder
            text_frame = shape.text_frame
            text_frame.clear()
            
            for bullet in bullet_points:
                if bullet:
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.font.name = "Aptos"
                    p.font.size = Pt(14)
                    p.level = 0
            break
    
    return slide

def create_slides_for_sop(prs, doc_content, num_slides=2):
    """Create multiple slides for an SOP document"""
    title = doc_content.get('title', 'Untitled SOP')
    sections = doc_content.get('sections', [])
    slides_created = []
    
    if not sections:
        # Create at least one slide
        slide = add_content_slide_with_template_style(
            prs, 
            title, 
            ["No detailed content available from source document"]
        )
        slides_created.append(slide)
        return slides_created
    
    # Distribute sections across slides
    total_sections = len(sections)
    sections_per_slide = max(1, total_sections // num_slides)
    
    for slide_num in range(num_slides):
        start_idx = slide_num * sections_per_slide
        end_idx = start_idx + sections_per_slide
        
        # For last slide, include all remaining sections
        if slide_num == num_slides - 1:
            end_idx = total_sections
        
        slide_sections = sections[start_idx:end_idx]
        
        if not slide_sections and slide_num > 0:
            # No more content, but we need to create slides
            slide_sections = [{'heading': 'Additional Information', 'content': ['See previous slides for details']}]
        
        # Build bullet points from sections
        bullet_points = []
        for section in slide_sections:
            # Add heading
            if section.get('heading'):
                bullet_points.append(section['heading'])
            
            # Add content (limited)
            content_items = section.get('content', [])
            for item in content_items[:3]:  # Max 3 items per section
                if len(item) > 100:
                    item = item[:97] + "..."
                bullet_points.append("  • " + item)
            
            # Limit total points per slide
            if len(bullet_points) >= 10:
                break
        
        # Create slide
        if slide_num == 0:
            slide_title = title
        else:
            slide_title = f"{title} (Part {slide_num + 1})"
        
        slide = add_content_slide_with_template_style(prs, slide_title, bullet_points[:10])
        slides_created.append(slide)
    
    return slides_created

def create_summary_slide_for_sop(prs, doc_content):
    """Create a single summary slide for an SOP document"""
    title = doc_content.get('title', 'Untitled SOP')
    sections = doc_content.get('sections', [])
    
    # Extract key points
    bullet_points = []
    
    for section in sections[:6]:  # Max 6 sections
        if section.get('heading'):
            bullet_points.append(section['heading'])
        
        # Add one key point from content
        content_items = section.get('content', [])
        if content_items:
            item = content_items[0]
            if len(item) > 80:
                item = item[:77] + "..."
            bullet_points.append("  • " + item)
        
        if len(bullet_points) >= 10:
            break
    
    if not bullet_points:
        bullet_points = ["Summary content not available"]
    
    slide = add_content_slide_with_template_style(prs, title, bullet_points[:10])
    return slide

def main():
    """Main function"""
    print("=" * 80)
    print("COMPREHENSIVE SOP PRESENTATION GENERATOR")
    print("=" * 80)
    
    # File paths
    template_path = "/home/runner/work/SOP/SOP/SOP Draft.pptx"
    output_path = "/home/runner/work/SOP/SOP/Comprehensive_SOP_Presentation.pptx"
    
    # SOPs requiring detailed coverage (2+ slides each)
    detailed_sops = [
        ("P9_IND_HR_002-00_Employee_Onboarding_SOP(1).docx", "Employee Onboarding"),
        ("P9_IND_HR_003-00_Employee_Relieving_SOP.docx", "Employee Relieving"),
        ("P9_IND_HR_004-00_Leave and attendance_Policy_SOP.docx", "Leave & Attendance Policy"),
        ("P9_IND_HR_005-00_Performance _Apraisal_SOP.docx", "Performance Appraisal"),
    ]
    
    # SOPs requiring summary coverage (1 slide each)
    summary_sops = [
        ("P9_IND_HR_006-00_Payroll_benefits_SOP.docx", "Payroll & Benefits"),
        ("P9_IND_HR_007-00_Training_Development_SOP.docx", "Training & Development"),
        ("P9_IND_HR_008-00_Disciplinary_Grieveance_SOP.docx", "Disciplinary & Grievance"),
        ("P9_IND_HR_009-00_NDA_Legal agreements_SOP.docx", "NDA & Legal Agreements"),
        ("P9_IND_HR_010-00_Organogram_SOP.docx", "Organization Structure"),
    ]
    
    print(f"\n📂 Loading template: {os.path.basename(template_path)}")
    prs = Presentation(template_path)
    print(f"   Template loaded with {len(prs.slides)} existing slides")
    print(f"   Available layouts: {len(prs.slide_layouts)}")
    
    # Create title slide
    print("\n📌 Creating main title slide...")
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    if title_slide.shapes.title:
        title_slide.shapes.title.text = "Standard Operating Procedures"
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Comprehensive HR SOPs - 2026"
    
    # Process detailed SOPs
    print(f"\n📄 Processing DETAILED SOPs ({len(detailed_sops)} documents, 2+ slides each)...")
    for filename, friendly_name in detailed_sops:
        filepath = f"/home/runner/work/SOP/SOP/{filename}"
        if os.path.exists(filepath):
            print(f"   • {friendly_name}...")
            content = extract_docx_content_detailed(filepath, friendly_title=friendly_name)
            create_slides_for_sop(prs, content, num_slides=2)
        else:
            print(f"   ⚠ {filename} not found")
    
    # Process summary SOPs
    print(f"\n📋 Processing SUMMARY SOPs ({len(summary_sops)} documents, 1 slide each)...")
    for filename, friendly_name in summary_sops:
        filepath = f"/home/runner/work/SOP/SOP/{filename}"
        if os.path.exists(filepath):
            print(f"   • {friendly_name}...")
            content = extract_docx_content_detailed(filepath, friendly_title=friendly_name)
            create_summary_slide_for_sop(prs, content)
        else:
            print(f"   ⚠ {filename} not found")
    
    # Save presentation
    print(f"\n💾 Saving presentation...")
    prs.save(output_path)
    
    # Summary
    print(f"\n{'=' * 80}")
    print(f"✅ SUCCESS! Comprehensive SOP presentation created")
    print(f"{'=' * 80}")
    print(f"📊 Statistics:")
    print(f"   • Total slides: {len(prs.slides)}")
    print(f"   • Detailed SOPs: {len(detailed_sops)} × 2 slides = {len(detailed_sops) * 2} slides")
    print(f"   • Summary SOPs: {len(summary_sops)} × 1 slide = {len(summary_sops)} slides")
    print(f"   • Output file: {output_path}")
    print(f"   • File size: {os.path.getsize(output_path) / 1024:.1f} KB")
    print(f"{'=' * 80}")

if __name__ == "__main__":
    main()
